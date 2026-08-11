#!/usr/bin/env python3
"""
Build Kobo_Pipeline_Simulator.pptx from the simulator.

The narration text is read directly out of Kobo_Pipeline_Simulator.html so the
deck always matches the on-screen voice-over. Run:

    python3 build_kobo_presentation.py
"""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
HTML = HERE / "Kobo_Pipeline_Simulator.html"
OUT = HERE / "Kobo_Pipeline_Simulator.pptx"

# 16:9 canvas
EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)

# palette lifted from the simulator's CSS variables
BG = RGBColor(0x07, 0x11, 0x1F)
PANEL = RGBColor(0x14, 0x2C, 0x4D)
PANEL2 = RGBColor(0x1B, 0x37, 0x60)
EDGE = RGBColor(0x27, 0x50, 0x7F)
BLUE = RGBColor(0x2B, 0x6C, 0xB2)
ACCENT = RGBColor(0x57, 0xA8, 0xE8)
GREEN = RGBColor(0x2F, 0xA3, 0x6B)
GREEN_INK = RGBColor(0x9F, 0xE3, 0xC1)
AMBER = RGBColor(0xE0, 0xA6, 0x3A)
RED = RGBColor(0xE0, 0x60, 0x5F)
INK = RGBColor(0xEA, 0xF1, 0xFA)
MUTED = RGBColor(0x93, 0xAB, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def load_narration():
    text = HTML.read_text()
    js = re.search(r"<script>(.*)</script>", text, re.S).group(1)
    narration = json.loads(re.search(r"var NARRATION = (\{.*?\n  \});", js, re.S).group(1))
    order = [el for el, _ in re.findall(r'\{el:"(\w+)",\s*steps:(\d+)', js)]
    return narration, order


NARR, ORDER = load_narration()


def lines(scene):
    return [item["t"] for item in NARR.get(scene, [])]


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    r._element.addprevious(r._element)  # keep behind
    return s


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, italic) in para:
            run = p.add_run()
            run.text = t
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = "Segoe UI"
    return tb


def R(t, size=14, color=INK, bold=False, italic=False):
    return (t, size, color, bold, italic)


def pill(slide, x, y, w, label, fill, line, txt_color):
    b = box(slide, x, y, w, Pt(24), fill=fill, line=line, line_w=1.0)
    tf = b.text_frame
    tf.word_wrap = False
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = txt_color
    run.font.name = "Segoe UI"
    return b


MX = int(0.62 * EMU)  # left margin


def header(slide, kicker, title, part_label=None, accent=ACCENT):
    # top accent rule
    box(slide, MX, int(0.55 * EMU), int(0.9 * EMU), Pt(4), fill=accent, radius=False)
    text(slide, MX, int(0.66 * EMU), SW - 2 * MX, int(0.4 * EMU),
         [[R(kicker.upper(), 12, accent, True, False)]])
    text(slide, MX, int(0.95 * EMU), SW - 2 * MX, int(0.9 * EMU),
         [[R(title, 30, WHITE, True, False)]])
    if part_label:
        pill(slide, SW - MX - int(2.7 * EMU), int(0.62 * EMU), int(2.7 * EMU),
             part_label, PANEL2, EDGE, ACCENT)


def footer(slide, idx):
    text(slide, MX, SH - int(0.5 * EMU), int(8 * EMU), int(0.35 * EMU),
         [[R("Kobo Tool Creator — pipeline walkthrough", 9, MUTED, False, False)]])
    text(slide, SW - MX - int(2 * EMU), SH - int(0.5 * EMU), int(2 * EMU), int(0.35 * EMU),
         [[R(str(idx), 9, MUTED, False, False)]], align=PP_ALIGN.RIGHT)


# rolling slide index for footer
_counter = {"n": 0}


def numbered(slide):
    _counter["n"] += 1
    footer(slide, _counter["n"])


# ----------------------------------------------------------------- title slide
def title_slide():
    s = slide()
    box(s, MX, int(2.35 * EMU), int(1.0 * EMU), Pt(5), fill=ACCENT, radius=False)
    text(s, MX, int(2.5 * EMU), SW - 2 * MX, int(1.6 * EMU),
         [[R("Kobo Tool Creator", 46, WHITE, True, False)],
          [R("From messy source data to five live Kobo forms", 22, ACCENT, False, False)]],
         space_after=8)
    text(s, MX, int(4.2 * EMU), int(9.5 * EMU), int(1.4 * EMU),
         [[R("An automated Google Apps Script pipeline that syncs the mentee and mentor "
             "databases, generates shared choice lists, then builds, validates and deploys "
             "every registered form to KoboToolbox.", 15, MUTED, False, False)]],
         line_spacing=1.2)
    # chips
    chips = ["1 trigger", "2 databases", "14 generated sheets", "5 Kobo tools", "1 EU server"]
    cx = MX
    for c in chips:
        w = int((0.42 + 0.092 * len(c)) * EMU)
        pill(s, cx, int(5.85 * EMU), w, c, PANEL, EDGE, GREEN_INK)
        cx += w + int(0.15 * EMU)
    numbered(s)


# ----------------------------------------------------------------- section divider
def divider(kicker, title, subtitle, part_label, accent=ACCENT):
    s = slide()
    box(s, MX, int(2.7 * EMU), int(1.0 * EMU), Pt(5), fill=accent, radius=False)
    text(s, MX, int(2.0 * EMU), SW - 2 * MX, int(0.5 * EMU),
         [[R(kicker.upper(), 13, accent, True, False)]])
    text(s, MX, int(2.85 * EMU), SW - 2 * MX, int(1.2 * EMU),
         [[R(title, 38, WHITE, True, False)]])
    text(s, MX, int(4.15 * EMU), int(10.5 * EMU), int(1.2 * EMU),
         [[R(subtitle, 17, MUTED, False, False)]], line_spacing=1.2)
    if part_label:
        pill(s, MX, int(1.35 * EMU), int(3.0 * EMU), part_label, PANEL2, EDGE, accent)
    numbered(s)


# ----------------------------------------------------------------- content: narration steps
def steps_slide(kicker, title, scene, part_label, accent=ACCENT, intro=None,
                two_col_threshold=5):
    s = slide()
    header(s, kicker, title, part_label, accent)
    items = lines(scene)
    top = int(2.05 * EMU)
    avail_h = SH - top - int(0.65 * EMU)
    if intro:
        text(s, MX, int(1.95 * EMU), SW - 2 * MX, int(0.5 * EMU),
             [[R(intro, 13, MUTED, False, True)]])
        top = int(2.45 * EMU)
        avail_h = SH - top - int(0.65 * EMU)

    two_col = len(items) > two_col_threshold
    if two_col:
        col_w = (SW - 2 * MX - int(0.4 * EMU)) // 2
        half = (len(items) + 1) // 2
        cols = [items[:half], items[half:]]
        xs = [MX, MX + col_w + int(0.4 * EMU)]
    else:
        col_w = SW - 2 * MX
        cols = [items]
        xs = [MX]

    for ci, col_items in enumerate(cols):
        row_h = avail_h // max(len(col_items), 1)
        row_h = min(row_h, int(1.25 * EMU))
        y = top
        start_num = 1 + (0 if ci == 0 else (len(items) + 1) // 2)
        for j, line in enumerate(col_items):
            n = start_num + j
            card = box(s, xs[ci], y, col_w, row_h - int(0.12 * EMU),
                       fill=PANEL, line=EDGE, line_w=1.0)
            # number badge
            bsz = Pt(26)
            badge = box(s, xs[ci] + Emu(int(0.12 * EMU)),
                        y + (row_h - int(0.12 * EMU)) // 2 - Emu(int(bsz / 2)),
                        bsz, bsz, fill=PANEL2, line=accent, line_w=1.0)
            btf = badge.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.margin_top = 0
            btf.margin_bottom = 0
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = str(n)
            br.font.size = Pt(12)
            br.font.bold = True
            br.font.color.rgb = accent
            br.font.name = "Segoe UI"
            # narration text
            text(s, xs[ci] + Emu(int(0.62 * EMU)),
                 y + Emu(int(0.02 * EMU)),
                 col_w - int(0.75 * EMU), row_h - int(0.12 * EMU),
                 [[R(line, 12.5, INK, False, False)]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
            y += row_h
    numbered(s)


# ----------------------------------------------------------------- tool summary slide
def tool_slide(order_no, title, tool_id, asset, fn, bullets, accent=GREEN):
    s = slide()
    header(s, f"Tool {order_no} of 5", title, f"{tool_id}", accent)
    # left: creation bullets
    lx = MX
    lw = int(7.0 * EMU)
    box(s, lx, int(2.1 * EMU), lw, int(4.4 * EMU), fill=PANEL, line=EDGE, line_w=1.0)
    text(s, lx + Emu(int(0.25 * EMU)), int(2.25 * EMU), lw - int(0.5 * EMU), int(0.4 * EMU),
         [[R("CREATION PROCESS", 12, accent, True, False)]])
    paras = []
    for b in bullets:
        paras.append([R("•  ", 14, accent, True, False), R(b, 13.5, INK, False, False)])
    text(s, lx + Emu(int(0.25 * EMU)), int(2.75 * EMU), lw - int(0.5 * EMU), int(3.6 * EMU),
         paras, space_after=9, line_spacing=1.12)

    # right: deploy facts
    rx = MX + lw + int(0.4 * EMU)
    rw = SW - MX - rx
    box(s, rx, int(2.1 * EMU), rw, int(4.4 * EMU), fill=PANEL2, line=EDGE, line_w=1.0)
    text(s, rx + Emu(int(0.22 * EMU)), int(2.25 * EMU), rw - int(0.44 * EMU), int(0.4 * EMU),
         [[R("VALIDATE · UPLOAD · DEPLOY", 12, accent, True, False)]])
    text(s, rx + Emu(int(0.22 * EMU)), int(2.8 * EMU), rw - int(0.44 * EMU), int(2.2 * EMU),
         [[R("Builder", 11, MUTED, False, False)],
          [R(fn + "()", 12.5, ACCENT, False, False)],
          [R("Kobo asset UID", 11, MUTED, False, False)],
          [R(asset, 12.5, GREEN_INK, False, False)],
          [R("Server", 11, MUTED, False, False)],
          [R("eu.kobotoolbox.org", 12.5, INK, False, False)]],
         space_after=6, line_spacing=1.05)
    pill(s, rx + Emu(int(0.22 * EMU)), int(5.7 * EMU), int(1.9 * EMU),
         "LIVE ✓", RGBColor(0x1D, 0x4F, 0x2E), GREEN, GREEN_INK)
    numbered(s)


# ----------------------------------------------------------------- five-tool overview
def registry_overview():
    s = slide()
    header(s, "Registry", "Five registered tools, built in this exact order", "buildRegisteredKoboTools_()")
    tools = [
        ("1", "EmONC Curriculum Tracking", "emonc_ctf", "aJaBJKDs7pCRMi8zm3BXze"),
        ("2", "Newborn Curriculum Tracking", "newborn_ctf", "a488FNw8rSGKWdJqpYfpny"),
        ("3", "MoH Skills Assessment", "moh_sac", "aR4bTSJFw3Tnev6o77S3Sg"),
        ("4", "Newborn Knowledge Assessment", "newborn_ka", "aFRcSLKi7wUvdrQ7js5Vbd"),
        ("5", "EmONC Knowledge Assessment", "emonc_ka", "auZqsTBpQMBnoDbMshmnrH"),
    ]
    gap = int(0.28 * EMU)
    cw = (SW - 2 * MX - 4 * gap) // 5
    y = int(2.4 * EMU)
    ch = int(3.7 * EMU)
    for i, (n, name, tid, asset) in enumerate(tools):
        x = MX + i * (cw + gap)
        box(s, x, y, cw, ch, fill=PANEL, line=EDGE, line_w=1.0)
        badge = box(s, x + cw // 2 - Emu(int(0.28 * EMU)), y + int(0.28 * EMU),
                    Pt(40), Pt(40), fill=PANEL2, line=ACCENT, line_w=1.25)
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = n
        br.font.size = Pt(18)
        br.font.bold = True
        br.font.color.rgb = ACCENT
        br.font.name = "Segoe UI"
        text(s, x + Emu(int(0.1 * EMU)), y + int(1.1 * EMU), cw - int(0.2 * EMU), int(1.7 * EMU),
             [[R(name, 13, WHITE, True, False)]], align=PP_ALIGN.CENTER, line_spacing=1.05)
        text(s, x + Emu(int(0.08 * EMU)), y + ch - int(1.15 * EMU), cw - int(0.16 * EMU), int(1.0 * EMU),
             [[R(tid, 10.5, ACCENT, False, False)],
              [R(asset, 8.5, GREEN_INK, False, False)]],
             align=PP_ALIGN.CENTER, space_after=3, line_spacing=1.0)
    text(s, MX, int(6.35 * EMU), SW - 2 * MX, int(0.5 * EMU),
         [[R("Each builder writes three tabs — survey, choices, settings — from the same "
             "generated sheets, then deploys in place to its existing Kobo asset.", 12.5, MUTED, False, True)]],
         align=PP_ALIGN.CENTER)
    numbered(s)


# ----------------------------------------------------------------- closing
def closing_slide():
    s = slide()
    box(s, MX, int(2.5 * EMU), int(1.0 * EMU), Pt(5), fill=GREEN, radius=False)
    text(s, MX, int(2.65 * EMU), SW - 2 * MX, int(1.0 * EMU),
         [[R("One trigger. Five live forms.", 36, WHITE, True, False)]])
    facts = [
        "Synced 2 hand-typed databases into clean local sheets",
        "Generated 14 shared choice / survey / logic sheets",
        "Built and validated 5 forms the way Kobo would",
        "Deployed 5 Kobo projects on eu.kobotoolbox.org",
        "Status = ok · 5 built · 5 deployed",
    ]
    paras = [[R("✓  ", 15, GREEN, True, False), R(f, 15, INK, False, False)] for f in facts]
    text(s, MX, int(3.9 * EMU), int(9.5 * EMU), int(2.6 * EMU), paras, space_after=8, line_spacing=1.1)
    text(s, MX, int(6.55 * EMU), SW - 2 * MX, int(0.5 * EMU),
         [[R("Interactive, narrated version: Kobo_Pipeline_Simulator.html", 12.5, ACCENT, False, False)]])
    numbered(s)


# ============================================================== assemble deck
title_slide()

# Part 1 — the process
divider("Part 1", "The Process",
        "One pipeline: messy source data → cleaned sheets → shared lists → validated forms → live on Kobo",
        "Part 1 · Process")
steps_slide("Step 0 · Trigger", "A weekly trigger reads the raw databases", "p1_trigger", "Part 1 · Process",
            intro="The mentee and mentor databases are typed by hand, so the data arrives messy.")
steps_slide("Steps 1–2 · Sync", "Clean the data into local sheets", "p1_sync", "Part 1 · Process",
            intro="Builders never read the raw databases — only these cleaned local copies.")
steps_slide("Step 3 · Generate", "Build the shared choice and survey lists", "p1_generate", "Part 1 · Process",
            intro="kobocreator turns the clean sheets into the lists every form reads.")
steps_slide("Step 4–5 · Build & Validate", "Write each form, then check it the way Kobo will", "p1_build", "Part 1 · Process")
steps_slide("Step 5 · Deploy", "Publish to Kobo and open the live form", "p1_deploy", "Part 1 · Process",
            intro="Only forms that passed validation are uploaded.")

# Part 2 — EmONC deep dive
divider("Part 2", "EmONC Curriculum Tracking",
        "Deep dive: watch one form assemble, then upload and deploy — Jane Wanjiku / Kirinyaga",
        "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Open form", "Prepare the three empty tabs", "p2_open", "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Choices", "Write the choices tab, row by row", "p2_choices", "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Survey", "Write the survey tab and the relevance cascade", "p2_survey", "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Outcome", "The finished form spreadsheet", "p2_outcome", "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Validate & Upload", "Check the form, export xlsx and import to Kobo", "p2_upload", "Part 2 · emonc_ctf", accent=GREEN)
steps_slide("Deploy", "Activate the new version — Jane is selectable", "p2_deploy", "Part 2 · emonc_ctf", accent=GREEN)

# Part 3 — the other tools
divider("Part 3", "The other four Kobo tools",
        "Each tool repeats the same safe pattern, in registry order",
        "Part 3 · Other tools", accent=AMBER)
registry_overview()
tool_slide("2", "Newborn Curriculum Tracking Form", "newborn_ctf", "a488FNw8rSGKWdJqpYfpny",
           "createNewbornCurriculumTrackingForm",
           ["Open the saved Newborn CTF spreadsheet",
            "Read the generated newborn survey, facility and mentee lists",
            "Write choices first: counties, facilities, active newborn mentees, modules",
            "Write survey + settings, dropping questions whose choice list is missing",
            "Validate, export xlsx, import and deploy the new version"])
tool_slide("3", "MoH Skills Assessment Checklist", "moh_sac", "aR4bTSJFw3Tnev6o77S3Sg",
           "createMoHSkillsAssessmentChecklist",
           ["Combine three programmes: IFM, Newborn and MENTORS survey rows",
            "Derive counties and facilities from All Facilities List (Choices)",
            "Kirinyaga and Kwale can never go missing — geography is data-derived",
            "Programme answer routes to the right facility, mentee and skills sections",
            "Validate, export, upload and deploy in place"])
tool_slide("4", "Newborn Knowledge Assessment", "newborn_ka", "aFRcSLKi7wUvdrQ7js5Vbd",
           "createNewbornKnowledgeAssessment",
           ["Questions from the Newborn Question Bank, or the built-in set",
            "Facilities from the generated Newborn Facilities List (Choices)",
            "Write questions, answers, constraints and score calculations",
            "Validate answer choices and calculations",
            "Export, upload and deploy the scoring form"])
tool_slide("5", "MoH Mentee EmONC Knowledge Assessment", "emonc_ka", "auZqsTBpQMBnoDbMshmnrH",
           "createEmONCKnowledgeAssessment",
           ["Questions from the EmONC Question Bank, or the built-in set",
            "Build the 16-county facility cascade from EmONC Facilities List",
            "Write assessment questions, relevance and scoring",
            "Validate county facility lists and assessment choices",
            "Export, upload and deploy — all five tools now live"])

# Part 4 — full pipeline
divider("Part 4", "The entire pipeline in one run",
        "One refreshAllKoboTools() run: prepare shared data once, build all five, then deploy",
        "Part 4 · Entire pipeline", accent=ACCENT)
steps_slide("Preparation", "Shared work that runs once", "p4_start", "Part 4 · Entire pipeline",
            intro="Lock, verify config and token, sync both databases, then generate all outputs.")
steps_slide("Build queue", "Build all five, in order", "p4_build_all", "Part 4 · Entire pipeline",
            intro="Each builder gets an isolated result; one failure never hides the others.")
steps_slide("Deploy queue", "Import and activate all five", "p4_deploy_all", "Part 4 · Entire pipeline",
            intro="Only successful builds are deployed — a stale sheet can never go live.")

closing_slide()

prs.save(str(OUT))
print(f"Wrote {OUT} with {len(prs.slides._sldIdLst)} slides")
